import json
import os
from typing import Any, Dict, Optional

import requests

try:
    from .config import ENABLE_CRAWL_ENRICHMENT, MAX_PITCH_DECK_BYTES
    from .website_utils import crawl_website
except ImportError:
    from config import ENABLE_CRAWL_ENRICHMENT, MAX_PITCH_DECK_BYTES
    from website_utils import crawl_website


OPENAI_API_URL = "https://api.openai.com/v1/chat/completions"
AI_INSIGHTS_MODEL = os.getenv("OPENAI_AI_INSIGHTS_MODEL", "gpt-4o")

SYSTEM_PROMPT = """You are a senior Venture Capital Partner evaluating startup investment opportunities.

Your job is to analyze founders, startups, and pitch decks from an investor's perspective.

Provide objective, professional, and actionable feedback.

Focus on:

* Founder quality
* Founder-market fit
* Startup attractiveness
* Market opportunity
* Business model
* Competitive advantage
* Traction
* Team strength
* Financial readiness
* Fundraising readiness
* Overall investment attractiveness

Return ONLY valid JSON."""


AI_INSIGHTS_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "properties": {
        "founder_analysis": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "overall_score": {"type": "integer"},
                "founder_market_fit": {"type": "integer"},
                "leadership_score": {"type": "integer"},
                "investor_confidence_score": {"type": "integer"},
                "strengths": {"type": "array", "items": {"type": "string"}},
                "weaknesses": {"type": "array", "items": {"type": "string"}},
                "recommendations": {"type": "array", "items": {"type": "string"}},
            },
            "required": [
                "overall_score",
                "founder_market_fit",
                "leadership_score",
                "investor_confidence_score",
                "strengths",
                "weaknesses",
                "recommendations",
            ],
        },
        "startup_analysis": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "overall_score": {"type": "integer"},
                "market_score": {"type": "integer"},
                "business_model_score": {"type": "integer"},
                "industry_score": {"type": "integer"},
                "fundraising_score": {"type": "integer"},
                "strengths": {"type": "array", "items": {"type": "string"}},
                "risks": {"type": "array", "items": {"type": "string"}},
                "recommendations": {"type": "array", "items": {"type": "string"}},
            },
            "required": [
                "overall_score",
                "market_score",
                "business_model_score",
                "industry_score",
                "fundraising_score",
                "strengths",
                "risks",
                "recommendations",
            ],
        },
        "pitch_deck_analysis": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "overall_score": {"type": "integer"},
                "investment_readiness": {"type": "string", "enum": ["High", "Medium", "Low"]},
                "recommendation": {"type": "string", "enum": ["Strong Buy", "Consider", "Monitor", "Pass"]},
                "executive_summary": {"type": "string"},
                "category_scores": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "problem": {"type": "integer"},
                        "solution": {"type": "integer"},
                        "market": {"type": "integer"},
                        "business_model": {"type": "integer"},
                        "traction": {"type": "integer"},
                        "team": {"type": "integer"},
                        "financials": {"type": "integer"},
                        "fundraising_readiness": {"type": "integer"},
                    },
                    "required": [
                        "problem",
                        "solution",
                        "market",
                        "business_model",
                        "traction",
                        "team",
                        "financials",
                        "fundraising_readiness",
                    ],
                },
                "strengths": {"type": "array", "items": {"type": "string"}},
                "weaknesses": {"type": "array", "items": {"type": "string"}},
                "investor_concerns": {"type": "array", "items": {"type": "string"}},
                "missing_information": {"type": "array", "items": {"type": "string"}},
                "improvement_recommendations": {"type": "array", "items": {"type": "string"}},
                "slide_analysis": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "additionalProperties": False,
                        "properties": {
                            "slide_number": {"type": "integer"},
                            "title": {"type": "string"},
                            "score": {"type": "integer"},
                            "strengths": {"type": "array", "items": {"type": "string"}},
                            "issues": {"type": "array", "items": {"type": "string"}},
                            "recommendations": {"type": "array", "items": {"type": "string"}},
                        },
                        "required": [
                            "slide_number",
                            "title",
                            "score",
                            "strengths",
                            "issues",
                            "recommendations",
                        ],
                    },
                },
            },
            "required": [
                "overall_score",
                "investment_readiness",
                "recommendation",
                "executive_summary",
                "category_scores",
                "strengths",
                "weaknesses",
                "investor_concerns",
                "missing_information",
                "improvement_recommendations",
                "slide_analysis",
            ],
        },
    },
    "required": ["founder_analysis", "startup_analysis", "pitch_deck_analysis"],
}


def extract_pitch_deck_text(file_bytes: bytes, filename: str, content_type: str = "") -> str:
    if len(file_bytes) > MAX_PITCH_DECK_BYTES:
        raise ValueError("Pitch deck exceeds the maximum file size of 10MB.")

    normalized_name = (filename or "").lower()
    if content_type == "application/pdf" or normalized_name.endswith(".pdf"):
        return extract_pdf_text(file_bytes)
    if (
        content_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or normalized_name.endswith(".pptx")
    ):
        return extract_pptx_text(file_bytes)
    raise ValueError("Only PDF and PPTX pitch deck files are accepted.")


def extract_pdf_text(file_bytes: bytes) -> str:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF pitch deck extraction.") from exc

    document = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        pages = []
        for index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"Slide {index}\n{text}")
        return "\n\n".join(pages)
    finally:
        document.close()


def extract_pptx_text(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX pitch deck extraction.") from exc

    from io import BytesIO

    presentation = Presentation(BytesIO(file_bytes))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text.strip())
        if parts:
            slides.append(f"Slide {index}\n" + "\n".join(parts))
    return "\n\n".join(slides)


def get_website_content(startup_data: Dict[str, Any]) -> str:
    website_url = startup_data.get("website_url") or startup_data.get("websiteUrl") or ""
    if not website_url or not ENABLE_CRAWL_ENRICHMENT:
        return ""
    return crawl_website(website_url)[:8000]


def build_user_prompt(
    founder_data: Dict[str, Any],
    startup_data: Dict[str, Any],
    website_content: str = "",
    pitch_deck_text: str = "",
) -> str:
    payload = {
        "founder_profile": founder_data,
        "startup_information": startup_data,
        "website_content": website_content,
        "pitch_deck_extracted_text": pitch_deck_text,
        "scoring_guidelines": {
            "90-100": "Investor Ready",
            "80-89": "Strong Candidate",
            "70-79": "Promising but Needs Improvement",
            "60-69": "Early Stage",
            "below_60": "Significant Gaps",
            "investment_readiness": ["High", "Medium", "Low"],
            "recommendation": ["Strong Buy", "Consider", "Monitor", "Pass"],
        },
    }
    return (
        "Analyze this startup for fundraising readiness. "
        "If pitch deck text is empty, analyze founder and startup information only and make the pitch deck section reflect missing deck context. "
        "Return JSON in the required schema.\n\n"
        f"{json.dumps(payload, ensure_ascii=False)}"
    )


def analyze_startup(
    founder_data: Dict[str, Any],
    startup_data: Dict[str, Any],
    pitch_deck_text: str = "",
    website_content: Optional[str] = None,
) -> Dict[str, Any]:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is not configured.")

    website_context = website_content if website_content is not None else get_website_content(startup_data)
    payload = {
        "model": AI_INSIGHTS_MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": build_user_prompt(founder_data, startup_data, website_context, pitch_deck_text[:30000])},
        ],
        "temperature": 0.2,
        "response_format": {
            "type": "json_schema",
            "json_schema": {
                "name": "ai_insights_analysis",
                "strict": True,
                "schema": AI_INSIGHTS_SCHEMA,
            },
        },
    }

    response = requests.post(
        OPENAI_API_URL,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        json=payload,
        timeout=60,
    )
    if response.status_code >= 400:
        raise RuntimeError(f"OpenAI analysis failed: {response.status_code} {response.text[:500]}")

    data = response.json()
    content = data["choices"][0]["message"]["content"]
    return json.loads(content)
